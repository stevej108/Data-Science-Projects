# ==========================================================
# Parser - PDF
# ==========================================================
from helpers.extraction import normalize_text, build_chunk, extract_ppt_text, extract_ppt_notes
from core.state import SearchState
from core.logger import get_logger
from core.config import EXCEL_ROWS_PER_CHUNK

logger = get_logger(__name__)

import fitz

def parse_pdf(document):

    logger.info(
        f"Parsing PDF | "
        f"{document['file_name']}"
    )

    chunks = []

    try:

        logger.info(
            "Opening PDF: %s",
            document["file_name"],
        )

        pdf = fitz.open(document["path"])

        logger.info(
            "Loaded %d pages.",
            len(pdf),
        )

        for page_num, page in enumerate(pdf, start=1):

            text = normalize_text(page.get_text())

            if not text:
                continue

            chunks.append(

                build_chunk(

                    document=document,

                    chunk_number=page_num,

                    section_title=f"Page {page_num}",

                    text=text,

                    metadata={
                        "page_number": page_num
                    }

                )

            )

        pdf.close()

        logger.info(
            f"Extracted {len(chunks)} pages"
        )

    except Exception:

        logger.exception(
            f"Failed parsing PDF | "
            f"{document['file_name']}"
        )

    return chunks



from docx import Document

def parse_word(document):

    logger.info(
        f"Parsing Word | "
        f"{document['file_name']} | "
        f"{document['document_id']}"
    )

    chunks = []

    try:

        logger.info(
            "Opening Word document: %s",
            document["file_name"],
        )

        doc = Document(document["path"])

        logger.info(
            "Loaded Word document."
        )

        heading = "Document"

        chunk_number = 1

        text_buffer = []

        for para in doc.paragraphs:

            text = normalize_text(para.text)

            if not text:
                continue

            # New heading starts a new chunk
            if para.style.name.startswith("Heading"):

                if text_buffer:

                    chunks.append(
                        build_chunk(
                            document=document,
                            chunk_number=chunk_number,
                            section_title=heading,
                            text="\n".join(text_buffer),
                            metadata={}
                        )
                    )

                    chunk_number += 1
                    text_buffer = []

                heading = text

            else:

                text_buffer.append(text)

        if text_buffer:

            chunks.append(
                build_chunk(
                    document=document,
                    chunk_number=chunk_number,
                    section_title=heading,
                    text="\n".join(text_buffer),
                    metadata={}
                )
            )

        logger.info(
            f"Extracted {len(chunks)} sections | "
            f"{document['file_name']}"
        )

    except Exception:

        logger.exception(
            f"Failed parsing Word | "
            f"{document['file_name']}"
        )

    return chunks

from openpyxl import load_workbook




def parse_excel(document):

    logger.info(
        "Parsing Excel | %s",
        document["file_name"],
    )

    chunks = []

    try:

        logger.info(
            "Opening workbook: %s",
            document["file_name"],
        )

        wb = load_workbook(
            document["path"],
            read_only=True,
            data_only=True,
        )

        logger.info(
            "Workbook contains %d worksheets.",
            len(wb.worksheets),
        )

        chunk_number = 1

        for sheet in wb.worksheets:

            logger.info(
                "Processing worksheet '%s' (rows=%d cols=%d)",
                sheet.title,
                sheet.max_row,
                sheet.max_column,
            )

            batch = []
            start_row = 1

            for row_idx, row in enumerate(
                sheet.iter_rows(values_only=True),
                start=1,
            ):

                values = [
                    str(v).strip()
                    for v in row
                    if v not in (None, "")
                ]

                if values:
                    batch.append(" | ".join(values))

                if len(batch) >= EXCEL_ROWS_PER_CHUNK:

                    text = normalize_text("\n".join(batch))

                    if text:

                        chunks.append(

                            build_chunk(

                                document=document,

                                chunk_number=chunk_number,

                                section_title=f"{sheet.title} (Rows {start_row}-{row_idx})",

                                text=text,

                                metadata={
                                    "worksheet": sheet.title,
                                    "row_start": start_row,
                                    "row_end": row_idx,
                                },

                            )

                        )

                        chunk_number += 1

                    batch = []
                    start_row = row_idx + 1

            # Remaining rows
            if batch:

                text = normalize_text("\n".join(batch))

                if text:

                    chunks.append(

                        build_chunk(

                            document=document,

                            chunk_number=chunk_number,

                            section_title=f"{sheet.title} (Rows {start_row}-{row_idx})",

                            text=text,

                            metadata={
                                "worksheet": sheet.title,
                                "row_start": start_row,
                                "row_end": row_idx,
                            },

                        )

                    )

                    chunk_number += 1

        wb.close()

        logger.info(
            "Extracted %d Excel chunks.",
            len(chunks),
        )

    except Exception:

        logger.exception(
            "Failed parsing Excel | %s",
            document["file_name"],
        )

    return chunks

# ==========================================================
# Parser - PowerPoint
# ==========================================================

from pptx import Presentation

def parse_powerpoint(document):

    logger.info(
        f"Parsing PowerPoint | "
        f"{document['file_name']} | "
        f"{document['document_id']}"
    )

    chunks = []

    try:

        logger.info(
            "Opening PowerPoint: %s",
            document["file_name"],
        )

        prs = Presentation(document["path"])

        logger.info(
            "Loaded %d slides.",
            len(prs.slides),
        )

        for slide_num, slide in enumerate(prs.slides, start=1):

            # ----------------------------
            # Extract title
            # ----------------------------
            title = ""

            if slide.shapes.title:
                title = normalize_text(
                    slide.shapes.title.text
                )

            # ----------------------------
            # Extract content
            # ----------------------------
            body = extract_ppt_text(slide)
            notes = extract_ppt_notes(slide)

            body = normalize_text(body)
            notes = normalize_text(notes)

            combined = "\n".join(
                x for x in [body, notes] if x
            )

            # ----------------------------
            # Build metadata
            # ----------------------------
            metadata = {

    "slide_number": slide_num,

    "has_notes": bool(notes),

    "shape_count": len(slide.shapes),

    "layout": slide.slide_layout.name
}

            chunk = build_chunk(

                document=document,

                chunk_number=slide_num,

                section_title=title,

                text=combined,

                metadata=metadata

            )

            chunks.append(chunk)

        logger.info(
            f"Extracted {len(chunks)} slides | "
            f"{document['file_name']}"
        )

    except Exception as e:

        logger.exception(
            f"Failed parsing PowerPoint | "
            f"{document['file_name']} | "
            f"{document['document_id']}"
        )

    return chunks




# ==========================================================
# Parser Registry
# ==========================================================

PARSERS = {
    ".ppt": parse_powerpoint,
    ".pptx": parse_powerpoint,
    ".pdf": parse_pdf,
    ".docx": parse_word,
    ".xlsx": parse_excel,
    ".xlsm": parse_excel,
}